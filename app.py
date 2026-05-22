# EC-AI Institutional Portfolio Intelligence v1
# Synthetic institutional banking portfolio + executive dashboard

import io
import os
import re
from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple

import numpy as np
import pandas as pd
import streamlit as st
import plotly.graph_objects as go
import plotly.express as px

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

# Optional exports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

# -----------------------------
# Page / theme
# -----------------------------
st.set_page_config(page_title="EC-AI Institutional Portfolio", layout="wide")

CONSULTING_PALETTE = [
    "#0B1F3B", "#2A6F97", "#2F855A", "#B7791F", "#9B2C2C",
    "#4A5568", "#718096", "#A0AEC0", "#76B7B2", "#E15759"
]

st.markdown(
    """
<style>
html, body, [class*="css"] { font-size: 16px; }
h1 { font-size: 38px !important; margin-bottom: 0.15rem; }
h2 { font-size: 26px !important; margin-top: 1.2rem; }
h3 { font-size: 20px !important; margin-top: 1.0rem; }
p, li { font-size: 16px; line-height: 1.5; }
.ec-kicker { color:#4B5563; font-size:18px; margin-bottom:4px; }
.ec-subtle { color:#6B7280; font-size:15px; }
.metric-card {
  border:1px solid rgba(17,24,39,0.10); border-radius:16px; padding:16px 18px;
  background:#fff; box-shadow:0 1px 2px rgba(17,24,39,0.04); min-height:118px;
}
.metric-label { color:#6B7280; font-size:13px; font-weight:700; text-transform:uppercase; letter-spacing:.03em; }
.metric-value { color:#111827; font-size:28px; font-weight:900; margin-top:6px; }
.metric-note { color:#374151; font-size:13px; margin-top:6px; }
.comment-box { border:1.5px dashed #C8CCD0; border-radius:12px; padding:16px 18px; background:#FAFBFC; }
</style>
""",
    unsafe_allow_html=True,
)

# -----------------------------
# Formatting
# -----------------------------
def fmt_money(x: float) -> str:
    try:
        x = float(x)
    except Exception:
        return "—"
    sign = "-" if x < 0 else ""
    x = abs(x)
    if x >= 1_000_000_000:
        return f"{sign}${x/1_000_000_000:.2f}B"
    if x >= 1_000_000:
        return f"{sign}${x/1_000_000:.1f}M"
    if x >= 1_000:
        return f"{sign}${x/1_000:.1f}K"
    return f"{sign}${x:,.0f}"


def fmt_pct(x: float, digits: int = 1) -> str:
    try:
        return f"{float(x)*100:.{digits}f}%"
    except Exception:
        return "—"


def fmt_bps(x: float) -> str:
    try:
        return f"{float(x)*10000:.0f} bps"
    except Exception:
        return "—"


def safe_div(a, b):
    try:
        return float(a) / float(b) if float(b) != 0 else np.nan
    except Exception:
        return np.nan


def apply_theme(fig: go.Figure, title: Optional[str] = None, height: int = 380, y_money: bool = False, y_pct: bool = False) -> go.Figure:
    fig.update_layout(
        template="plotly_white",
        title=dict(text=title or fig.layout.title.text, x=0.0, xanchor="left", font=dict(size=18, color="#111827")),
        height=height,
        margin=dict(l=46, r=26, t=58, b=48),
        font=dict(family="Inter, Arial, sans-serif", size=13, color="#111827"),
        colorway=CONSULTING_PALETTE,
        paper_bgcolor="white",
        plot_bgcolor="white",
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="left", x=0),
    )
    fig.update_xaxes(showline=True, linecolor="#374151", linewidth=1, ticks="outside", gridcolor="rgba(17,24,39,0.06)")
    fig.update_yaxes(showline=True, linecolor="#374151", linewidth=1, ticks="outside", gridcolor="rgba(17,24,39,0.06)", rangemode="tozero")
    if y_money:
        fig.update_yaxes(tickprefix="$", tickformat=",.2s")
    if y_pct:
        fig.update_yaxes(tickformat=".0%")
    return fig


def metric_card(label: str, value: str, note: str = ""):
    st.markdown(
        f"""
<div class="metric-card">
  <div class="metric-label">{label}</div>
  <div class="metric-value">{value}</div>
  <div class="metric-note">{note}</div>
</div>
""",
        unsafe_allow_html=True,
    )

# -----------------------------
# Synthetic institutional dataset generator
# -----------------------------
SECTORS = ["Technology", "Real Estate", "Utilities", "Manufacturing", "Healthcare", "Consumer", "Transport", "Energy", "Financial Institutions"]
COUNTRIES = ["Hong Kong", "Singapore", "Japan", "China", "Australia", "India", "Korea", "Taiwan", "Thailand", "Indonesia"]
PRODUCTS = ["Term Loan", "RCF", "Trade Finance", "FX", "Deposits", "Cash Management", "Bond", "Syndicated Loan"]
FACILITY_TYPES = ["Committed", "Uncommitted", "Bilateral", "Syndicated", "Working Capital"]
CURRENCIES = ["USD", "HKD", "JPY", "SGD", "AUD", "CNY"]
RMS = ["RM Alpha", "RM Bravo", "RM Charlie", "RM Delta", "RM Echo", "RM Foxtrot", "RM Gamma", "RM Helix"]
GROUP_PREFIX = ["Apex", "Pacific", "Orient", "Harbour", "Summit", "Nexus", "Crescent", "Phoenix", "Evergreen", "Vertex", "Golden", "Atlas"]
GROUP_SUFFIX = ["Holdings", "Group", "Industries", "Capital", "International", "Partners", "Logistics", "Energy", "Properties", "Technologies"]


def make_synthetic_portfolio(n_clients: int = 160, seed: int = 42) -> pd.DataFrame:
    rng = np.random.default_rng(seed)
    rows = []
    groups = [f"{rng.choice(GROUP_PREFIX)} {rng.choice(GROUP_SUFFIX)}" for _ in range(max(40, n_clients // 2))]

    for i in range(n_clients):
        group = rng.choice(groups)
        country = rng.choice(COUNTRIES, p=[.16,.13,.12,.13,.09,.09,.08,.08,.06,.06])
        sector = rng.choice(SECTORS)
        rm = rng.choice(RMS)
        product = rng.choice(PRODUCTS, p=[.18,.15,.12,.13,.18,.10,.06,.08])
        facility = rng.choice(FACILITY_TYPES)
        ccy = rng.choice(CURRENCIES, p=[.45,.17,.12,.11,.07,.08])
        risk_rating = int(np.clip(np.round(rng.normal(5.2, 1.8)), 1, 10))
        rel_year = int(rng.integers(2004, 2026))

        base = rng.lognormal(mean=18.7, sigma=0.85)  # institutional scale
        sector_mult = {
            "Real Estate": 1.45, "Financial Institutions": 1.65, "Technology": 1.25,
            "Utilities": 1.15, "Energy": 1.35, "Consumer": .85
        }.get(sector, 1.0)
        country_mult = {"Japan": 1.25, "China": 1.2, "Hong Kong": 1.15, "Singapore": 1.12}.get(country, 1.0)

        credit_limit = base * sector_mult * country_mult
        utilization = float(np.clip(rng.beta(3.0, 2.8), .05, .98))
        loan_outstanding = credit_limit * utilization if product not in ["Deposits", "Cash Management", "FX"] else credit_limit * rng.uniform(.02, .25)
        deposit_balance = credit_limit * rng.uniform(.05, 1.15) if product in ["Deposits", "Cash Management", "FX", "RCF"] else credit_limit * rng.uniform(.01, .45)
        rwa_density = float(np.clip(.28 + risk_rating * .055 + rng.normal(0, .08), .18, .95))
        rwa = loan_outstanding * rwa_density
        nim = float(np.clip(rng.normal(.018 + risk_rating*.0018, .005), .004, .055))
        fee_margin = float(np.clip(rng.normal(.006, .003), .001, .018))
        loan_interest_income = loan_outstanding * nim
        fee_income = credit_limit * fee_margin * rng.uniform(.25, 1.1)
        trade_finance_income = credit_limit * rng.uniform(.0005, .006) if product in ["Trade Finance", "Cash Management", "RCF"] else credit_limit * rng.uniform(0, .0015)
        fx_income = credit_limit * rng.uniform(.0003, .0045) if product in ["FX", "Cash Management", "Trade Finance"] else credit_limit * rng.uniform(0, .0012)
        revenue = loan_interest_income + fee_income + trade_finance_income + fx_income
        capital_charge = rwa * .11
        roe = safe_div(revenue * rng.uniform(.38, .62), capital_charge)
        wallet_size = revenue / rng.uniform(.08, .42)
        wallet_share = safe_div(revenue, wallet_size)
        expected_usage = credit_limit * float(np.clip(utilization + rng.normal(0.04, .11), .03, 1.0))
        pd_val = float(np.clip(.0015 * (risk_rating ** 1.65) * rng.uniform(.55, 1.6), .0005, .25))
        lgd = float(np.clip(rng.normal(.42, .13), .15, .75))
        stage = "Stage 1" if risk_rating <= 6 else ("Stage 2" if risk_rating <= 8 else "Stage 3")
        watchlist = bool(risk_rating >= 8 or rng.random() < .05)
        covenant = bool(risk_rating >= 8 and rng.random() < .28)
        esg = float(np.clip(rng.normal(52 + (sector in ["Energy", "Real Estate"])*12, 18), 5, 98))
        casa_ratio = float(np.clip(rng.beta(2.2, 3.2), .02, .92))
        funding_months = int(np.clip(rng.normal(18, 11), 1, 60))
        maturity_bucket = rng.choice(["0-3M", "3-6M", "6-12M", "1-3Y", "3Y+"], p=[.16,.17,.25,.28,.14])
        deposit_mix = rng.choice(["Operating Deposit", "Term Deposit", "Non-interest Bearing", "Investment Sweep"], p=[.32,.34,.22,.12])
        liquidity_profile = rng.choice(["Sticky", "Moderate", "Price Sensitive", "Volatile"], p=[.30,.38,.22,.10])
        cross_sell_score = float(np.clip(35 + 55*(1-wallet_share) + rng.normal(0,12), 0, 100))
        ai_growth_score = float(np.clip(20 + cross_sell_score*.45 + (roe < .12)*15 + rng.normal(0,10), 0, 100))
        relationship_strength = float(np.clip(65 + (2026-rel_year)*1.3 - risk_rating*2 + rng.normal(0,12), 0, 100))
        expansion_potential = float(np.clip((1-wallet_share)*80 + relationship_strength*.15 + rng.normal(0,8), 0, 100))
        momentum = float(np.clip(rng.normal(50 + ai_growth_score*.22 - risk_rating*1.5, 13), 0, 100))

        rows.append({
            "Client Name": f"Client {i+1:03d} - {group.split()[0]}",
            "Group Name": group,
            "Country": country,
            "Sector": sector,
            "RM": rm,
            "Booking Center": country,
            "Facility Type": facility,
            "Product Type": product,
            "Currency": ccy,
            "Risk Rating": risk_rating,
            "Relationship Start Year": rel_year,
            "Credit Limit": credit_limit,
            "Loan Outstanding": loan_outstanding,
            "Deposit Balance": deposit_balance,
            "RWA": rwa,
            "Revenue": revenue,
            "ROE": roe,
            "NIM": nim,
            "Utilization %": utilization,
            "Wallet Size": wallet_size,
            "Wallet Share": wallet_share,
            "Expected Usage": expected_usage,
            "Fee Income": fee_income,
            "Trade Finance Income": trade_finance_income,
            "FX Income": fx_income,
            "PD": pd_val,
            "LGD": lgd,
            "Stage": stage,
            "Watchlist Flag": watchlist,
            "Covenant Breach Flag": covenant,
            "ESG Risk Score": esg,
            "CASA Ratio": casa_ratio,
            "Funding Duration Months": funding_months,
            "Maturity Bucket": maturity_bucket,
            "Deposit Mix": deposit_mix,
            "Liquidity Profile": liquidity_profile,
            "Cross Sell Score": cross_sell_score,
            "AI Growth Score": ai_growth_score,
            "Relationship Strength": relationship_strength,
            "Expansion Potential": expansion_potential,
            "Client Momentum": momentum,
        })

    df = pd.DataFrame(rows)
    return df

# -----------------------------
# Data model and insights
# -----------------------------
@dataclass
class PortfolioSummary:
    revenue: float
    loans: float
    deposits: float
    rwa: float
    roe: float
    nim: float
    utilization: float
    wallet_share: float
    watchlist_count: int
    client_count: int


def summarize(df: pd.DataFrame) -> PortfolioSummary:
    revenue = df["Revenue"].sum()
    loans = df["Loan Outstanding"].sum()
    deposits = df["Deposit Balance"].sum()
    rwa = df["RWA"].sum()
    roe = safe_div((df["Revenue"].sum() * .50), (df["RWA"].sum() * .11))
    nim = safe_div((df["Loan Outstanding"] * df["NIM"]).sum(), df["Loan Outstanding"].sum())
    utilization = safe_div(df["Loan Outstanding"].sum(), df["Credit Limit"].sum())
    wallet_share = safe_div(df["Revenue"].sum(), df["Wallet Size"].sum())
    watchlist_count = int(df["Watchlist Flag"].sum())
    return PortfolioSummary(revenue, loans, deposits, rwa, roe, nim, utilization, wallet_share, watchlist_count, len(df))


def build_exec_insights(df: pd.DataFrame, s: PortfolioSummary) -> List[str]:
    top_country = df.groupby("Country")["Revenue"].sum().sort_values(ascending=False)
    top_sector = df.groupby("Sector")["Revenue"].sum().sort_values(ascending=False)
    low_roe = df[df["ROE"] < .10]
    high_growth = df.sort_values("AI Growth Score", ascending=False).head(5)
    watch = df[df["Watchlist Flag"]]
    deposit_mix = df.groupby("Deposit Mix")["Deposit Balance"].sum().sort_values(ascending=False)

    lines = [
        f"Portfolio covers **{s.client_count:,} institutional clients** with total revenue of **{fmt_money(s.revenue)}**, loans of **{fmt_money(s.loans)}**, and deposits of **{fmt_money(s.deposits)}**.",
        f"Overall portfolio ROE is **{fmt_pct(s.roe)}** with average NIM of **{fmt_bps(s.nim)}** and utilization of **{fmt_pct(s.utilization)}**.",
    ]
    if not top_country.empty:
        share = top_country.iloc[0] / s.revenue
        lines.append(f"Revenue concentration is led by **{top_country.index[0]}**, contributing **{fmt_money(top_country.iloc[0])}** or **{fmt_pct(share)}** of total revenue.")
    if not top_sector.empty:
        lines.append(f"Top sector is **{top_sector.index[0]}** with **{fmt_money(top_sector.iloc[0])}** revenue; monitor sector concentration against risk appetite.")
    if len(low_roe):
        lines.append(f"There are **{len(low_roe)} low-ROE relationships** below 10%; these are candidates for repricing, utilization improvement, or balance-sheet reduction.")
    if len(watch):
        lines.append(f"Watchlist exposure: **{len(watch)} clients** with combined loan outstanding of **{fmt_money(watch['Loan Outstanding'].sum())}**.")
    if not deposit_mix.empty:
        lines.append(f"Deposit franchise is led by **{deposit_mix.index[0]}** deposits at **{fmt_money(deposit_mix.iloc[0])}**; quality of funding should be checked, not just balance size.")
    if len(high_growth):
        names = ", ".join(high_growth["Client Name"].head(3).tolist())
        lines.append(f"Highest AI growth score names include **{names}** — strong candidates for wallet growth and cross-sell campaigns.")
    return lines

# -----------------------------
# Charts
# -----------------------------
def bar_chart(df: pd.DataFrame, x: str, y: str, title: str, topn: int = 10, y_money: bool = True) -> go.Figure:
    d = df.groupby(x, dropna=False)[y].sum().sort_values(ascending=False).head(topn).reset_index()
    fig = go.Figure(go.Bar(x=d[x].astype(str), y=d[y], text=d[y], texttemplate="%{text:,.2s}", textposition="outside", marker_color=CONSULTING_PALETTE[:len(d)]))
    return apply_theme(fig, title, height=380, y_money=y_money)


def heatmap_roe(df: pd.DataFrame) -> go.Figure:
    piv = df.pivot_table(index="Sector", columns="Country", values="ROE", aggfunc="mean")
    fig = px.imshow(piv, text_auto=".1%", aspect="auto", color_continuous_scale="RdYlGn", zmin=0, zmax=max(.25, float(np.nanpercentile(df["ROE"], 90))))
    fig.update_layout(title=dict(text="ROE Heatmap — Sector x Country", x=0, font=dict(size=18)), height=430, margin=dict(l=60, r=30, t=60, b=50), font=dict(size=12))
    return fig


def scatter_roe_risk(df: pd.DataFrame) -> go.Figure:
    fig = px.scatter(
        df,
        x="ROE",
        y="Risk Rating",
        size="Loan Outstanding",
        color="Sector",
        hover_name="Client Name",
        hover_data={"Country": True, "Revenue": ":,.0f", "Wallet Share": ":.1%", "Loan Outstanding": ":,.0f"},
    )
    fig.add_vline(x=.10, line_dash="dash", line_color="#9B2C2C", annotation_text="10% ROE floor")
    fig.update_yaxes(autorange="reversed")
    return apply_theme(fig, "Client Map — ROE vs Risk Rating", height=430, y_money=False)


def maturity_ladder(df: pd.DataFrame) -> go.Figure:
    order = ["0-3M", "3-6M", "6-12M", "1-3Y", "3Y+"]
    d = df.groupby("Maturity Bucket")["Loan Outstanding"].sum().reindex(order).fillna(0).reset_index()
    fig = go.Figure(go.Bar(x=d["Maturity Bucket"], y=d["Loan Outstanding"], marker_color=CONSULTING_PALETTE[1], text=d["Loan Outstanding"], texttemplate="%{text:,.2s}", textposition="outside"))
    return apply_theme(fig, "Maturity Ladder — Loan Outstanding", height=360, y_money=True)


def deposit_mix_donut(df: pd.DataFrame) -> go.Figure:
    d = df.groupby("Deposit Mix")["Deposit Balance"].sum().sort_values(ascending=False).reset_index()
    fig = go.Figure(go.Pie(labels=d["Deposit Mix"], values=d["Deposit Balance"], hole=.58, textinfo="label+percent", marker=dict(colors=CONSULTING_PALETTE)))
    fig.update_layout(title=dict(text="Deposit Mix", x=0, font=dict(size=18)), height=360, margin=dict(l=20, r=20, t=60, b=20), legend=dict(orientation="h", y=-.05))
    return fig


def wallet_scatter(df: pd.DataFrame) -> go.Figure:
    fig = px.scatter(
        df,
        x="Wallet Share",
        y="Expansion Potential",
        size="Revenue",
        color="Relationship Strength",
        hover_name="Client Name",
        hover_data={"Country": True, "Sector": True, "Revenue": ":,.0f", "AI Growth Score": ":.1f"},
        color_continuous_scale="Blues",
    )
    fig.add_vline(x=.25, line_dash="dash", line_color="#4A5568", annotation_text="Wallet share 25%")
    return apply_theme(fig, "Wallet Opportunity Map", height=430, y_pct=False)

# -----------------------------
# Ask AI
# -----------------------------
def get_api_key() -> Optional[str]:
    try:
        if hasattr(st, "secrets") and "OPENAI_API_KEY" in st.secrets:
            return st.secrets["OPENAI_API_KEY"]
    except Exception:
        pass
    return os.environ.get("OPENAI_API_KEY")


def build_ai_context(df: pd.DataFrame, summary_lines: List[str]) -> str:
    s = summarize(df)
    top_clients = df.sort_values("Revenue", ascending=False).head(10)[["Client Name", "Country", "Sector", "Revenue", "ROE", "Risk Rating", "Wallet Share", "AI Growth Score"]]
    low_roe = df.sort_values("ROE").head(10)[["Client Name", "Country", "Sector", "Revenue", "ROE", "Risk Rating", "Loan Outstanding"]]
    watch = df[df["Watchlist Flag"]].sort_values("Loan Outstanding", ascending=False).head(10)[["Client Name", "Country", "Sector", "Loan Outstanding", "Risk Rating", "PD", "LGD"]]
    return f"""
You are EC-AI Institutional Portfolio Intelligence. Answer using only these facts. Do not invent data.

PORTFOLIO TOTALS
Clients: {s.client_count}
Revenue: {fmt_money(s.revenue)}
Loans: {fmt_money(s.loans)}
Deposits: {fmt_money(s.deposits)}
RWA: {fmt_money(s.rwa)}
ROE: {fmt_pct(s.roe)}
NIM: {fmt_bps(s.nim)}
Utilization: {fmt_pct(s.utilization)}
Wallet Share: {fmt_pct(s.wallet_share)}
Watchlist count: {s.watchlist_count}

EXECUTIVE INSIGHTS
{chr(10).join('- ' + re.sub(r'[*`]', '', x) for x in summary_lines)}

TOP CLIENTS BY REVENUE
{top_clients.to_string(index=False)}

LOWEST ROE CLIENTS
{low_roe.to_string(index=False)}

WATCHLIST EXPOSURE
{watch.to_string(index=False)}
"""


def ask_ai(question: str, context: str) -> str:
    if not question.strip():
        return "Please enter a question."
    key = get_api_key()
    if not key:
        return "OpenAI API key not configured. Add OPENAI_API_KEY in Streamlit secrets to enable Ask AI."
    if OpenAI is None:
        return "OpenAI SDK not installed. Add openai to requirements.txt."
    try:
        client = OpenAI(api_key=key)
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            temperature=.2,
            max_tokens=450,
            messages=[
                {"role": "system", "content": context},
                {"role": "user", "content": f"Question: {question}\nAnswer with specific numbers and concise executive recommendations."},
            ],
        )
        return (resp.choices[0].message.content or "").strip()
    except Exception as e:
        return f"Ask AI error: {e}"

# -----------------------------
# Export helpers
# -----------------------------
def fig_to_png(fig: go.Figure) -> bytes:
    return fig.to_image(format="png", scale=2)


def build_ppt(deck_title: str, charts: List[Tuple[str, go.Figure, str]]) -> bytes:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed.")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = deck_title
    slide.placeholders[1].text = "Executive Pack — Institutional Portfolio Intelligence"
    blank = prs.slide_layouts[6]
    for title, fig, bullets in charts:
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(Inches(.6), Inches(.35), Inches(12), Inches(.6))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        img = io.BytesIO(fig_to_png(fig))
        slide.shapes.add_picture(img, Inches(.6), Inches(1.15), width=Inches(7.4))
        bx = slide.shapes.add_textbox(Inches(8.25), Inches(1.15), Inches(4.6), Inches(5.7))
        tf = bx.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = "Commentary"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(16)
        for line in bullets.split("\n")[:7]:
            pp = tf.add_paragraph()
            pp.text = line.strip("-• ")
            pp.font.size = Pt(13)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

# -----------------------------
# UI
# -----------------------------
st.title("EC-AI Institutional Portfolio Intelligence")
st.markdown("<div class='ec-kicker'>Corporate banking portfolio, liquidity, risk, and relationship economics — explained clearly.</div>", unsafe_allow_html=True)
st.markdown("<div class='ec-subtle'>Synthetic dataset mode lets you demo institutional analytics without exposing confidential bank data.</div>", unsafe_allow_html=True)
st.divider()

with st.sidebar:
    st.header("Portfolio Data")
    mode = st.radio("Data source", ["Synthetic institutional portfolio", "Upload CSV"], index=0)
    n_clients = st.slider("Synthetic clients", 50, 500, 180, step=10)
    seed = st.number_input("Random seed", min_value=1, max_value=9999, value=42, step=1)
    uploaded = st.file_uploader("Upload portfolio CSV", type=["csv"])
    st.caption("Expected fields: Revenue, Loan Outstanding, Deposit Balance, RWA, ROE, NIM, Country, Sector, Product Type, Risk Rating.")

if mode == "Upload CSV" and uploaded is not None:
    df = pd.read_csv(uploaded)
else:
    df = make_synthetic_portfolio(n_clients=n_clients, seed=int(seed))

required = ["Revenue", "Loan Outstanding", "Deposit Balance", "RWA", "ROE", "NIM", "Country", "Sector", "Product Type", "Risk Rating"]
missing = [c for c in required if c not in df.columns]
if missing:
    st.error(f"Missing required columns: {missing}")
    st.stop()

# Filters
with st.expander("Filters", expanded=True):
    c1, c2, c3, c4 = st.columns(4)
    with c1:
        countries = st.multiselect("Country", sorted(df["Country"].unique()), default=sorted(df["Country"].unique()))
    with c2:
        sectors = st.multiselect("Sector", sorted(df["Sector"].unique()), default=sorted(df["Sector"].unique()))
    with c3:
        products = st.multiselect("Product", sorted(df["Product Type"].unique()), default=sorted(df["Product Type"].unique()))
    with c4:
        max_rr = st.slider("Max Risk Rating", 1, 10, 10)

fdf = df[df["Country"].isin(countries) & df["Sector"].isin(sectors) & df["Product Type"].isin(products) & (df["Risk Rating"] <= max_rr)].copy()
if fdf.empty:
    st.warning("No records after filters.")
    st.stop()

s = summarize(fdf)
insights = build_exec_insights(fdf, s)

# KPI cards
st.subheader("Executive Portfolio Summary")
r1 = st.columns(4)
with r1[0]: metric_card("Revenue", fmt_money(s.revenue), f"{s.client_count:,} active relationships")
with r1[1]: metric_card("Loans", fmt_money(s.loans), f"Utilization {fmt_pct(s.utilization)}")
with r1[2]: metric_card("Deposits", fmt_money(s.deposits), "Funding franchise view")
with r1[3]: metric_card("RWA", fmt_money(s.rwa), "Capital consumption")
r2 = st.columns(4)
with r2[0]: metric_card("Portfolio ROE", fmt_pct(s.roe), "Revenue / capital proxy")
with r2[1]: metric_card("Avg NIM", fmt_bps(s.nim), "Weighted by loan balance")
with r2[2]: metric_card("Wallet Share", fmt_pct(s.wallet_share), "Revenue captured vs estimated wallet")
with r2[3]: metric_card("Watchlist", f"{s.watchlist_count}", "Clients flagged for monitoring")

st.divider()
st.subheader("Executive Insights")
for line in insights:
    st.markdown(f"- {line}")

st.divider()
st.subheader("Institutional Dashboard")

fig_country = bar_chart(fdf, "Country", "Revenue", "Revenue by Country", topn=10)
fig_sector = bar_chart(fdf, "Sector", "Loan Outstanding", "Loan Outstanding by Sector", topn=10)
fig_roe = heatmap_roe(fdf)
fig_risk = scatter_roe_risk(fdf)
fig_mat = maturity_ladder(fdf)
fig_dep = deposit_mix_donut(fdf)
fig_wallet = wallet_scatter(fdf)
fig_product = bar_chart(fdf, "Product Type", "Revenue", "Revenue by Product", topn=10)

row = st.columns(2, gap="large")
with row[0]: st.plotly_chart(fig_country, use_container_width=True, config={"displayModeBar": False})
with row[1]: st.plotly_chart(fig_sector, use_container_width=True, config={"displayModeBar": False})
row = st.columns(2, gap="large")
with row[0]: st.plotly_chart(fig_roe, use_container_width=True, config={"displayModeBar": False})
with row[1]: st.plotly_chart(fig_risk, use_container_width=True, config={"displayModeBar": False})
row = st.columns(2, gap="large")
with row[0]: st.plotly_chart(fig_mat, use_container_width=True, config={"displayModeBar": False})
with row[1]: st.plotly_chart(fig_dep, use_container_width=True, config={"displayModeBar": False})
row = st.columns(2, gap="large")
with row[0]: st.plotly_chart(fig_wallet, use_container_width=True, config={"displayModeBar": False})
with row[1]: st.plotly_chart(fig_product, use_container_width=True, config={"displayModeBar": False})

st.divider()
st.subheader("Priority Relationship Actions")

low_roe = fdf[fdf["ROE"] < .10].sort_values(["Loan Outstanding", "Revenue"], ascending=False).head(10)
growth = fdf.sort_values("AI Growth Score", ascending=False).head(10)
watch = fdf[fdf["Watchlist Flag"]].sort_values("Loan Outstanding", ascending=False).head(10)

c1, c2, c3 = st.columns(3, gap="large")
with c1:
    st.markdown("#### Reprice / Improve ROE")
    st.dataframe(low_roe[["Client Name", "Country", "Sector", "Revenue", "ROE", "Loan Outstanding", "Risk Rating"]], use_container_width=True, hide_index=True)
with c2:
    st.markdown("#### Cross-sell / Grow Wallet")
    st.dataframe(growth[["Client Name", "Country", "Sector", "Revenue", "Wallet Share", "AI Growth Score", "Expansion Potential"]], use_container_width=True, hide_index=True)
with c3:
    st.markdown("#### Risk Monitor")
    st.dataframe(watch[["Client Name", "Country", "Sector", "Loan Outstanding", "Risk Rating", "PD", "LGD", "Stage"]], use_container_width=True, hide_index=True)

st.divider()
st.subheader("Ask AI (Portfolio Q&A)")
st.caption("Ask: Which clients should we reprice? Where is concentration risk? Which segment has wallet growth potential?")
context = build_ai_context(fdf, insights)
if "portfolio_ai_history" not in st.session_state:
    st.session_state.portfolio_ai_history = []
q_col, b_col = st.columns([.82, .18])
with q_col:
    q = st.text_input("Ask EC-AI…", placeholder="Which relationships should we prioritize next quarter?")
with b_col:
    clicked = st.button("Ask", use_container_width=True)
if clicked and q.strip():
    with st.spinner("Thinking…"):
        ans = ask_ai(q, context)
    st.session_state.portfolio_ai_history.insert(0, (q, ans))
    st.session_state.portfolio_ai_history = st.session_state.portfolio_ai_history[:5]
for qq, aa in st.session_state.portfolio_ai_history[:3]:
    st.markdown(f"**Q:** {qq}")
    st.markdown(aa)

st.divider()
st.subheader("Export / Download")

csv = fdf.to_csv(index=False).encode("utf-8")
st.download_button("Download Filtered Portfolio CSV", csv, file_name="ecai_institutional_portfolio_dataset_v1.csv", mime="text/csv")

charts = [
    ("Revenue by Country", fig_country, "Country concentration shows where portfolio revenue is generated. Compare revenue concentration with capital and risk concentration."),
    ("Loan Outstanding by Sector", fig_sector, "Sector exposure shows balance-sheet deployment. Use this to identify concentration and diversification needs."),
    ("ROE Heatmap", fig_roe, "ROE heatmap highlights high-return and low-return country-sector pockets. Low-return cells should trigger repricing or balance-sheet review."),
    ("Client Map — ROE vs Risk", fig_risk, "Client-level view separates attractive relationships from high-risk, low-return names."),
    ("Maturity Ladder", fig_mat, "Maturity ladder supports liquidity and refinancing planning."),
    ("Wallet Opportunity Map", fig_wallet, "Low wallet share and high expansion potential names are the commercial growth targets."),
]

if st.button("Generate Executive Pack (PPTX)", use_container_width=True):
    try:
        pptx = build_ppt("EC-AI Institutional Portfolio Intelligence v1", charts)
        st.download_button("Download PPTX", pptx, file_name="ecai_institutional_portfolio_v1.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
    except Exception as e:
        st.error(f"PPTX generation failed: {e}")

with st.expander("Data preview", expanded=False):
    st.dataframe(fdf, use_container_width=True, hide_index=True)
